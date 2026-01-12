import os
import subprocess
import re
import sys
from PIL import Image
import io
from pptx import Presentation
from pptx.util import Inches
from pptx.shapes.freeform import FreeformBuilder

# Ensure potrace is installed: `apt-get install potrace`

def trace_bitmap_to_svg(image_path: str, output_svg: str):
    """
    Uses potrace to trace a bitmap (BMP/PBM) to SVG.
    """
    # 1. Convert to BMP (Potrace prefers BMP or PBM)
    with Image.open(image_path) as img:
        # Convert to black and white (1-bit)
        # Thresholding
        gray = img.convert('L')
        bw = gray.point(lambda x: 0 if x < 128 else 255, '1')
        bmp_path = image_path + ".bmp"
        bw.save(bmp_path)

    # 2. Run Potrace
    try:
        # -b svg: Backend SVG
        # --alphamax 0: Turn off curve smoothing optimization for stricter architecture lines? Maybe default is better.
        # -k 0.5: Black level
        cmd = ["potrace", bmp_path, "-s", "-o", output_svg]
        subprocess.run(cmd, check=True)
        print(f"Vectorized to {output_svg}")
    except FileNotFoundError:
        print("Error: 'potrace' command not found. Please install it.")
        return False
    except subprocess.CalledProcessError as e:
        print(f"Potrace failed: {e}")
        return False
    finally:
        if os.path.exists(bmp_path):
            os.remove(bmp_path)
    return True

def parse_svg_path(svg_file):
    """
    Very naive SVG path parser.
    Extracts 'd' attributes from <path> tags.
    Returns a list of path strings.
    """
    paths = []
    with open(svg_file, 'r') as f:
        content = f.read()
        # Regex to find d="..."
        # This is robust enough for Potrace output which is usually simple 
        # but huge single paths.
        matches = re.findall(r'<path[^>]*\sd="([^"]+)"', content)
        paths.extend(matches)
    return paths

def apply_svg_path_to_shape(slide, path_str, width_scale=0.01, height_scale=0.01):
    """
    Converts SVG path data to PPTX Freeform shape.
    PPTX FreeformBuilder:
    - .add_line_segments(vertices, close=True)
    - But SVG has Curves (C, Q, S). python-pptx freeform support is primarily line segments?
    - Actually FreeformBuilder has no 'curve' method in older versions. 
    - Ideally we approximate curves or use .convert_to_shape() if building manually.
    
    Constraint: python-pptx FreeformBuilder mainly supports 'add_line_segments'.
    We might need to linearize curves or just hope Potrace output (which uses beziers) 
    can be flattened or accept lines.
    
    Let's assume we approximate everything with lines for this prototype.
    """
    
    # SVG Commands regex
    # M x y
    # L x y
    # C x1 y1 x2 y2 x y
    # Z
    
    # Potrace uses absolute coordinates usually?
    tokens = re.findall(r'([a-zA-Z])|([-+]?\d*\.?\d+(?:[eE][-+]?\d+)?)', path_str)
    
    # Flatten tokens
    commands = []
    for t in tokens:
        if t[0]: commands.append(t[0])
        if t[1]: commands.append(float(t[1]))
        
    cursor_x, cursor_y = 0, 0
    start_x, start_y = 0, 0
    
    # Scale factor for PPTX (EMU or Inches)
    # Potrace coordinates are pixel-based.
    SCALE = 7200 # Arbitrary scaling factor
    
    # We need to construct points for add_line_segments
    # [(x, y), (x, y), ...]
    
    vertices = []
    
    idx = 0
    while idx < len(commands):
        cmd = commands[idx]
        idx += 1
        
        if cmd == 'M':
            x = commands[idx] * SCALE
            y = commands[idx+1] * SCALE
            idx += 2
            cursor_x, cursor_y = x, y
            start_x, start_y = x, y
            vertices.append((x, y))
            
        elif cmd == 'L':
            x = commands[idx] * SCALE
            y = commands[idx+1] * SCALE
            idx += 2
            cursor_x, cursor_y = x, y
            vertices.append((x, y))
            
        elif cmd == 'C':
            # Cubic Bezier: x1 y1 x2 y2 x y
            # We ignore control points for this crude prototype and just draw line to end
            # To do it right, we'd sample the bezier.
            x1 = commands[idx] * SCALE
            y1 = commands[idx+1] * SCALE
            x2 = commands[idx+2] * SCALE
            y2 = commands[idx+3] * SCALE
            x = commands[idx+4] * SCALE
            y = commands[idx+5] * SCALE
            idx += 6
            cursor_x, cursor_y = x, y
            vertices.append((x, y))

        elif cmd == 'Z':
            # Close path
            vertices.append((start_x, start_y))
            
            # Flush shape
            if len(vertices) > 2:
                # Build shape
                try:
                    shapes = slide.shapes
                    # build_freeform available?
                    freeform_builder = shapes.build_freeform(vertices[0][0], vertices[0][1])
                    freeform_builder.add_line_segments(vertices[1:], close=False)
                    shape = freeform_builder.convert_to_shape()
                    # Apply style
                    fill = shape.fill
                    fill.solid()
                    fill.fore_color.rgb = 0x000000 # Black
                    line = shape.line
                    line.width = 0
                except Exception as e:
                    print(f"Failed to build shape: {e}")
            
            vertices = []

# Main Test Logic
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python research_vectorize.py <image_path>")
        sys.exit(1)
        
    src_img = sys.argv[1]
    svg_out = "temp_trace.svg"
    pptx_out = "vector_test.pptx"
    
    print(f"Tracing {src_img}...")
    if trace_bitmap_to_svg(src_img, svg_out):
        print(f"Parsing {svg_out}...")
        paths = parse_svg_path(svg_out)
        print(f"Found {len(paths)} paths.")
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        for p in paths:
            apply_svg_path_to_shape(slide, p)
            
        prs.save(pptx_out)
        print(f"Saved to {pptx_out}")
    else:
        print("Tracing failed.")
