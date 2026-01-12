from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def analyze_pptx(path):
    print(f"Analyzing: {path}")
    try:
        prs = Presentation(path)
    except Exception as e:
        print(f"Failed to open PPTX: {e}")
        return

    print(f"Total Slides: {len(prs.slides)}")
    print("-" * 30)

    for i, slide in enumerate(prs.slides):
        if i >= 3: break # Analyze first 3 slides
        print(f"Slide {i+1}:")
        print(f"  Background: {slide.background.fill.type if slide.background else 'None'}")
        
        shapes = slide.shapes
        print(f"  Total Shapes: {len(shapes)}")
        
        for shape in shapes:
            shape_type = shape.shape_type
            name = shape.name
            
            info = f"    - [{shape_type}] '{name}'"
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                w, h = shape.width, shape.height
                info += f" (Image: {w}x{h})"
                # Try to guess format or transparency? Difficult via pptx directly without blobbing.
                
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                info += f" (Group with {len(shape.shapes)} items)"
                
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                info += f" (AutoShape)"
                if shape.has_text_frame:
                    info += f" [Text: {shape.text[:20]}...]"
            
            elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                if shape.has_text_frame:
                    info += f" [Text: {shape.text[:20]}...]"
                    
            print(info)
        print("-" * 30)

if __name__ == "__main__":
    analyze_pptx("/Users/bangdoll/Library/Mobile Documents/iCloud~md~obsidian/Documents/AI筆記/youtube_summary/Docs/Awakening_Blueprint.pptx")
