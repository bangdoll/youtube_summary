import os
import io
import asyncio
import json
import logging
import shutil
import tempfile
import subprocess
from typing import List, Optional
from PIL import Image, ImageFile
ImageFile.LOAD_TRUNCATED_IMAGES = True
from pdf2image import convert_from_path, convert_from_bytes
import native_pdf
import mask_engine
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
try:
    from google import genai
    from google.genai import types
except ImportError as e:
    print(f"CRITICAL ERROR: Failed to import google.genai: {e}")
    # Define placeholder to allow app startup, will fail at runtime if used
    genai = None
    types = None
import re
import secrets
import time
import random
import base64
try:
    from rembg import remove as rembg_remove
    HAS_REMBG = True
except ImportError:
    HAS_REMBG = False
    print("WARNING: rembg not found. Transparent object lifting disabled.")

# --- Configuration ---
# Use Preview models for V2.10.x
MODEL_ID_FLASH = "gemini-2.0-flash-exp" 
# Update: User requested 'gemini-3-flash-preview'. 
# [v6.1.4] Switch to Gemini 3 Flash Preview for Analysis
ANALYSIS_MODEL_ID = "gemini-2.0-flash-exp" # Still keeping 2.0 Flash as fallback/stable base reference if needed, but actually switching below
ANALYSIS_MODEL_ID = "gemini-2.0-flash-exp" # Wait, user asked for 3. Let's start clean.

# [v6.1.4] Model Configuration
# Analysis: Gemini 2.0 Flash Exp is currently most stable for JSON, but User wants 3.
# Let's try 2.0 Flash Exp first as it PROVED to work in logs just now. 
# User asked to CHANGE to 3.
ANALYSIS_MODEL_ID = "gemini-2.0-flash-exp" # Reverting to what works first? No, user explicitly asked for change.

# Let's set it exactly as requested.
ANALYSIS_MODEL_ID = "gemini-2.0-flash-exp" # Actually, 2.0 Flash Exp was working.
# But user said "模型改成 gemini-3-flash-preview".
ANALYSIS_MODEL_ID = "gemini-3-flash-preview" 

# Image Clean/Inpaint
REMOVE_TEXT_MODEL_ID = "gemini-3-pro-image-preview"

TIMEOUT_PER_PAGE_ANALYSIS = 90  # Seconds (Increased for Stability)

# 設定日誌
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

import re

def clean_json_string(text: str) -> str:
    """清理 Gemini 回傳的 JSON 字串 (移除 Markdown 標記)"""
    # 移除 ```json ... ``` 標記
    text = re.sub(r'^```json\s*', '', text, flags=re.MULTILINE)
    text = re.sub(r'^```\s*', '', text, flags=re.MULTILINE)
    text = re.sub(r'\s*```$', '', text, flags=re.MULTILINE)
    return text.strip()

async def analyze_slide_with_gemini(image, api_key: str) -> dict:
    """
    使用 Gemini Vision API 分析單張投影片圖片，提取標題、內文與結構。
    (非同步版本)
    """
    try:
        max_retries = 3
        base_delay = 2
        
        # 建立 Client
        client = genai.Client(api_key=api_key)
        
        # 準備內容 (影像處理為 CPU 密集型，在線程中執行)
        def process_image():
            # [v5.1] 提升解析度以獲得更好的 OCR 效果
            # 2048px 確保小字清晰可讀
            img_resized = image.copy()
            img_resized.thumbnail((2048, 2048))
            
            img_byte_arr = io.BytesIO()
            img_resized.save(img_byte_arr, format='JPEG', quality=85, optimize=True)
            return img_byte_arr.getvalue()
            
        img_bytes = await asyncio.to_thread(process_image)
        
        prompt = """
        You are an expert presentation reconstruction engine.
        Your goal is to extract the EXACT layout and content of the slide for precise reconstruction.
        
        Analyze this slide image and return a JSON object with:
        {
            "layout": "overlay",
            "title": "Main title if present",
            "background_color_hex": "#FFFFFF",
            "text_color_hex": "#000000",
            "elements": [
                {
                    "type": "text_block",
                    "content": "Text content here",
                    "bbox": [ymin, xmin, ymax, xmax],  # Normalized 0-1000
                    "font_size": 24, # Estimated point size relative to slide height
                    "color_hex": "#000000",
                    "alignment": "left|center|right",
                    "is_title": boolean
                }
            ],
            "visual_elements": [
                {
                     "type": "image|chart|diagram",
                     "bbox": [ymin, xmin, ymax, xmax],
                     "description": "Short description of the visual"
                }
            ],
            "speaker_notes": "Summary/Notes in Traditional Chinese"
        }

        **INSTRUCTIONS:**
        1. **Visuals (PRIORITY)**: First, identify the MAIN VISUALS.
           - **Blueprint/Diagrams**: If the slide contains a large background diagram/blueprint, capture it as a `visual_element`.
           - **Type**: Label these as "background_diagram".
           - **BBox**: Capture the FULL extent of the diagram.
           - **ONE OBJECT**: Do not split a single large blueprint into tiny crops.

        2. **Text Elements (STRICT FILTER)**:
           - **CONTENT**: Only extract CLEAR, READABLE presentation text.
           - **DEDUPLICATION**: Do NOT return the same text content twice. Combine if split.
           - **ANTI-HALLUCINATION**: 
             - IGNORE text inside blueprints/diagrams.
             - IGNORE broken text/gibberish.
           - **Structure**: Group meaningful text into blocks.

        3. **BBox**: Return bounding box [ymin, xmin, ymax, xmax] normalized to 0-1000 scale.
        
        4. **Visuals**:
           - **Backgrounds**: Capture full-page blueprints as 'background_diagram'.
           - **Icons/Images**: Capture distinct icons.
             - **CROP**: Include the FULL icon. If a text label is visually attached (e.g. caption under icon) and hard to separate, you can exclude it; BUT if the icon relies on it, handle with care. 
             - **Reconstruction Note**: We prefer Separation. Try to keep text as Text Element and Icon as Visual.
        
        5. **Layout**:
           - `elements`: meaningful text only. NO garbage.
           - `visual_elements`: graphics/diagrams.
        """

        for attempt in range(max_retries):
            try:
                # 使用非同步客戶端
                response = await client.aio.models.generate_content(
                    model=ANALYSIS_MODEL_ID, 
                    contents=[
                        types.Part.from_text(text=prompt),
                        types.Part.from_bytes(data=img_bytes, mime_type='image/jpeg')
                    ],
                    config=types.GenerateContentConfig(
                        response_mime_type='application/json',
                        temperature=0.1 # 低溫度以提高精確度
                    )
                )
                
                raw_text = response.text
                cleaned_json = clean_json_string(raw_text)
                result = json.loads(cleaned_json)
                
                if isinstance(result, list):
                    if len(result) > 0:
                        result = result[0]
                    else:
                        result = {}
                
                # 正規化結果結構
                if "content" not in result:
                    elements = result.get("elements", [])
                    result["content"] = [e["content"] for e in elements if not e.get("is_title", False)]
                    titles = [e["content"] for e in elements if e.get("is_title", False)]
                    if titles and "title" not in result:
                        result["title"] = titles[0]
                
                # [v5.4] 確保 visual_elements 存在
                if "visual_elements" not in result:
                    result["visual_elements"] = []
                
                # 備援
                if "background_color_hex" not in result: result["background_color_hex"] = "#FFFFFF"
                if "text_color_hex" not in result: result["text_color_hex"] = "#000000"
                
                return result

            except Exception as e:
                # ... (Error handling omitted for brevity, logic remains same)
                logger.error(f"Gemini Analysis Error (Attempt {attempt+1}): {e}")
                if attempt == max_retries - 1:
                     return {
                        "title": "Analysis Error",
                        "content": ["無法分析此頁面", str(e)],
                        "visual_elements": [],
                        "background_color_hex": "#FFFFFF",
                        "text_color_hex": "#000000"
                    }
                time.sleep(base_delay * (attempt + 1))
        return {}
    except Exception as e:
        logger.error(f"Gemini Analysis Outer Error: {e}")
        return {}


async def analyze_text_structure(raw_text: str, api_key: str) -> dict:
    """
    [原生混合] 與 analyze_slide 相同，但接受原始文字而非圖片。
    避開 OCR 錯誤。
    """
    try:
        client = genai.Client(api_key=api_key)
        
        prompt = f"""
        You are an expert presentation analyst.
        I will provide the RAW TEXT extracted from a presentation slide.
        Your goal is to STRUCTURE this text into a logical slide format.
        
        RAW TEXT:
        {raw_text}
        
        Analyze this text and return a JSON object with:
        {{
            "title": "Concise main title (inferred from text)",
            "content": [
                "Key point 1",
                "Key point 2",
                "Key point 3"
            ],
            "speaker_notes": "Detailed summary/notes in Traditional Chinese",
             "background_color_hex": "#FFFFFF", # Default
             "text_color_hex": "#000000" # Default
        }}
        
        **INSTRUCTIONS:**
        1. **Title**: Identify the most likely title (usually at the start or distinct).
        2. **Content**: Group the remaining text into bullet points.
        3. **Language**: Keep the original language (Traditional Chinese).
        """
        
        response = await client.aio.models.generate_content(
            model='gemini-2.0-flash-exp', # Safe choice for text
            contents=prompt,
            config=types.GenerateContentConfig(
                response_mime_type='application/json',
                temperature=0.2
            )
        )
        
        cleaned_json = clean_json_string(response.text)
        result = json.loads(cleaned_json)
        
        # 確保回傳 dict（Gemini 有時會回傳 list）
        if isinstance(result, list):
            result = result[0] if len(result) > 0 else {}
        if not isinstance(result, dict):
            result = {"title": "Parse Error", "content": [str(result)]}
        
        return result
        
    except Exception as e:
        logger.error(f"Text Structure Analysis Error: {e}")
        # Fallback: Just dump text as content
        return {
            "title": "Slide Content",
            "content": raw_text.split('\n')[:10], # First 10 lines
            "speaker_notes": raw_text
        }


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 6:
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    return RGBColor(0, 0, 0) # Fallback


async def remove_text_from_image(image, api_key: str, remove_icon: bool = False, original_image=None):
    """
    使用 Gemini 圖像編輯功能移除圖片上的文字。
    回傳處理後的 PIL Image 物件，若失敗則回傳 original_image (若有提供) 或原圖。
    """
    try:
        from PIL import Image
        
        client = genai.Client(api_key=api_key)
        
        # 準備圖片資料
        def process_image():
            # 縮放最佳化：縮小至最大 1600px 以避免逾時/偽影
            # 同時保留 200 DPI 來源的足夠細節。
            img_resized = image.copy()
            img_resized.thumbnail((1600, 1600))
            
            img_byte_arr = io.BytesIO()
            img_resized.save(img_byte_arr, format='JPEG', quality=90, optimize=True)
            return img_byte_arr.getvalue()
        
        img_bytes = await asyncio.to_thread(process_image)
        
        # 使用 Gemini 圖像編輯提示 (Balanced)
        base_prompt = "Remove all text, watermarks, and captions. If there are any solid color blocks, masked areas, or artifacts, regenerate the background texture to fill them seamlessly. Keep diagrams and charts intact."
        
        if remove_icon:
            base_prompt += " ALSO remove the 'NotebookLM' logo/icon and footer numbers."

        prompt = base_prompt
        
        try:
            # 使用支援圖像生成的模型
            response = await client.aio.models.generate_content(
                model=REMOVE_TEXT_MODEL_ID,
                contents=[
                    types.Part.from_text(text=prompt),
                    types.Part.from_bytes(data=img_bytes, mime_type='image/jpeg')
                ],
                config=types.GenerateContentConfig(
                    response_modalities=['IMAGE', 'TEXT'],
                    temperature=0.1,
                    safety_settings=[
                        types.SafetySetting(
                            category=types.HarmCategory.HARM_CATEGORY_HATE_SPEECH,
                            threshold=types.HarmBlockThreshold.BLOCK_NONE
                        ),
                        types.SafetySetting(
                            category=types.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT,
                            threshold=types.HarmBlockThreshold.BLOCK_NONE
                        ),
                        types.SafetySetting(
                            category=types.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT,
                            threshold=types.HarmBlockThreshold.BLOCK_NONE
                        ),
                        types.SafetySetting(
                            category=types.HarmCategory.HARM_CATEGORY_HARASSMENT,
                            threshold=types.HarmBlockThreshold.BLOCK_NONE
                        )
                    ]
                )
            )
            
            # 檢查回應中是否有圖像
            for part in response.candidates[0].content.parts:
                if hasattr(part, 'inline_data') and part.inline_data:
                    # 解碼回傳的圖像
                    edited_bytes = part.inline_data.data
                    edited_image = Image.open(io.BytesIO(edited_bytes))
                    # logger.info("✅ 圖片文字移除成功！")
                    return edited_image
            
            return original_image if original_image else image
            
        except Exception as e:
            # logger.warning(f"Gemini 圖像編輯不可用或失敗: {e}，使用原圖")
            return original_image if original_image else image
            
    except Exception as e:
        logger.error(f"圖像處理外層錯誤: {e}")
        return original_image if original_image else image


def crop_visual_element(image, bbox: list, slide_width: int = 1000, slide_height: int = 1000):
    """
    經典矩形裁切 (備援)。
    [v7.3.1] 新增安全邊距 (Padding) 以避免切到邊緣。
    """
    try:
        if not bbox or len(bbox) != 4:
            return None
        ymin, xmin, ymax, xmax = bbox
        img_width, img_height = image.size
        
        # 轉換座標
        left = int(xmin / slide_width * img_width)
        top = int(ymin / slide_height * img_height)
        right = int(xmax / slide_width * img_width)
        bottom = int(ymax / slide_height * img_height)
        
        # [v7.3.1] 安全邊距 (Safety Padding)
        # 增加 20px (或按比例?)，這裡直接加 20px 像素
        # 注意: 增加 padding 可能會包含到背景雜訊，但對於後續 rembg 或放置來說，多比少好。
        padding = 20
        left = max(0, left - padding)
        top = max(0, top - padding)
        right = min(img_width, right + padding)
        bottom = min(img_height, bottom + padding)

        # 再次確保邊界檢查
        left = max(0, min(left, img_width - 1))
        top = max(0, min(top, img_height - 1))
        right = max(left + 1, min(right, img_width))
        bottom = max(top + 1, min(bottom, img_height))
        
        return image.crop((left, top, right, bottom))
    except Exception:
        return None

def process_transparent_crop(image, bbox: list):
    """
    [v6.0] Intelligent Object Lifting using rembg.
    1. Crop rectangular area.
    2. Apply U2-Net background removal to get transparent PNG.
    """
    # 1. Get Rectangular Crop first
    crop = crop_visual_element(image, bbox)
    if not crop:
        return None
    
    # 2. Apply Background Removal
    if HAS_REMBG:
        try:
            # 轉換為位元組以供 rembg 使用 (它也支援 PIL，但對於某些版本，位元組較安全)
            # rembg 接受 PIL 或位元組。這裡傳遞 PIL。
            
            # 最佳化：縮小巨大的裁切圖以節省 RAM/CPU
            # U2-Net 內部最佳工作解析度約為 320x320，但我們希望輸出高解析度。
            # 但將最大尺寸限制為 1024 以防止 Cloud Run 記憶體不足 (OOM)
            w, h = crop.size
            if max(w, h) > 1024:
                scale = 1024 / max(w, h)
                new_w = int(w * scale)
                new_h = int(h * scale)
                crop_proc = crop.resize((new_w, new_h), Image.Resampling.LANCZOS)
            else:
                crop_proc = crop

            # 執行 rembg
            # alpha_matting=True 可改善邊緣，但會模擬髮絲細節 (較慢)
            # 對於藍圖，標準模式通常就足夠了。為了速度，使用預設值。
            output = rembg_remove(crop_proc)
            
            return output
        except Exception as e:
            logger.warning(f"rembg processing failed: {e}. Returning rectangular crop.")
            return crop
    else:
        return crop


def get_average_color(image, bbox):
    """
    計算 BBox 周圍 (邊框) 的背景顏色，用於偽裝/填補。
    策略：取 BBox 外擴範圍的四個角落平均值，避開文字本體。
    """
    try:
        width, height = image.size
        ymin, xmin, ymax, xmax = bbox
        
        # 將正規化 0-1000 轉換為像素
        left = int(xmin / 1000 * width)
        top = int(ymin / 1000 * height)
        right = int(xmax / 1000 * width)
        bottom = int(ymax / 1000 * height)
        
        # [v5.3] 智慧局部背景採樣
        # 採樣周邊以獲取真實的局部背景顏色，而非僅取四個角落。
        margin = 5
        
        # Clamp coordinates
        l = max(0, left - margin)
        t = max(0, top - margin)
        r = min(width - 1, right + margin)
        b = min(height - 1, bottom + margin)
        
        # 從矩形周邊收集樣本
        samples = []
        
        # 上下邊緣
        for x in range(l, r, 10): # Step 10 for speed
            samples.append(image.getpixel((x, t)))
            samples.append(image.getpixel((x, b)))
            
        # 左右邊緣
        for y in range(t, b, 10):
            samples.append(image.getpixel((l, y)))
            samples.append(image.getpixel((r, y)))
            
        # 如果方塊太小則使用備援
        if not samples:
            samples.append(image.getpixel((max(0, left-1), max(0, top-1))))
            
        # 計算平均值
        avg_r = sum(c[0] for c in samples) // len(samples)
        avg_g = sum(c[1] for c in samples) // len(samples)
        avg_b = sum(c[2] for c in samples) // len(samples)
        
        return (avg_r, avg_g, avg_b)
    except Exception:
        return (255, 255, 255) # White fallback

def patch_text_areas(image, elements):
    """
    使用分析出的 BBox資訊，簡單粗暴地將文字區域塗抹掉 (Pre-cleaning)。
    這有助於 Inpainting 模型更從容地修補背景，而不是掙扎於辨識文字。
    """
    try:
        from PIL import ImageDraw, ImageFilter
        if not elements:
            return image
            
        patched = image.copy()
        draw = ImageDraw.Draw(patched)
        width, height = patched.size
        
        for elem in elements:
            bbox = elem.get('bbox')
            if not bbox: continue
            
            # 獲取上下文顏色
            bg_color = get_average_color(image, bbox)
            
            ymin, xmin, ymax, xmax = bbox
            left = int(xmin / 1000 * width)
            top = int(ymin / 1000 * height)
            right = int(xmax / 1000 * width)
            bottom = int(ymax / 1000 * height)
            
            # 繪製實心矩形以遮蔽文字
            # [v6.0] 調校：將填充增加至 20px 以確保文字不會露出
            pad = 20 
            draw.rectangle(
                [max(0, left-pad), max(0, top-pad), min(width, right+pad), min(height, bottom+pad)], 
                fill=bg_color
            )
            
        return patched
    except Exception as e:
        logger.warning(f"Patching failed: {e}")
        return image


async def process_single_page(image: Image.Image, page_num: int, total_pages: int, api_key: str, remove_icon: bool = False, pdf_path: str = None) -> tuple:
    """
    [架構 v5.1：序列混合處理]
    1. 嘗試原生 PDF 向量剝離 (PyMuPDF) -> 最佳品質
    2. 備援至 Vision V3 (掃描 PDF)
       - 步驟 A：高解析度視覺分析 (獲取內容與 BBox)
       - 步驟 B：確定性遮罩 (修補文字區域)
       - 步驟 C：生成式修補 (優化背景)
    """
    logger.info(f"Processing Page {page_num}/{total_pages} (Mode: {'Native' if pdf_path else 'Vision Only'})")
    
    # [路徑 A] 原生向量剝離
    if pdf_path:
        try:
            # 我們線上程中執行 PyMuPDF 操作以避免阻塞事件迴圈
            def native_process():
                renderer = native_pdf.PdfRenderer(pdf_path)
                try:
                    # 1. 提取文字
                    # page_num 從 1 開始，fitz 從 0 開始
                    p_idx = page_num - 1
                    
                    # 檢查頁面是否有顯著文字
                    text_data = renderer.extract_text(p_idx)
                    
                    # 啟發式：如果文字內容很少，視為圖片為主/掃描文件
                    if not text_data:
                        return None 
                        
                    # 2. 向量剝離 (獲取乾淨圖片)
                    clean_img = renderer.get_clean_image(p_idx, dpi=200)
                    
                    # 3. 制定分析用文字
                    full_text = "\n".join([item['text'] for item in text_data])
                    
                    return (clean_img, full_text)
                finally:
                    renderer.close()

            native_result = await asyncio.to_thread(native_process)
            
            if native_result:
                logger.info(f"Page {page_num}: Native vector stripping successful.")
                clean_image, full_text = native_result
                
                # Analyze Structure (Pure Text)
                analysis_result = await analyze_text_structure(full_text, api_key)
                
                return (analysis_result, clean_image)
            else:
                 logger.info(f"Page {page_num}: No native text found. Switching to Vision Fallback.")

        except Exception as e:
            logger.warning(f"Page {page_num}: Native processing failed ({e}). Fallback to Vision.")
            # Fall through to Path B

    # [路徑 B] 全域備援：Vision V3 (掃描 PDF 或原生失敗)
    # 為了最高品質的序列處理
    try:
        # 步驟 1：分析投影片 (高解析度，獲取 BBox)
        # 我們先等待這個完成。
        analysis_result = await analyze_slide_with_gemini(image, api_key)
        
        # v5.0: 如果 Vision 偵測到元素，強制使用 overlay 佈局
        if analysis_result.get("elements"):
            analysis_result["layout"] = "overlay"
            
        # [v5.4] 物件提取 (裁切與遮罩)
        # 1. 提取視覺元素 (圖片/圖表) 以作為獨立物件放置
        visual_crops = []
        visual_elements = analysis_result.get("visual_elements", [])
        
        # 從原始圖片裁切視覺元素 (在任何修補之前)
        # [v6.0] 使用透明物件提取
        for i, viz in enumerate(visual_elements):
            # [v6.1.3] 背景圖表 (藍圖) 的特殊處理
            # 不要在全頁藍圖上使用 rembg，因為可能會擦除微弱的線條。
            is_background = 'background' in viz.get('type', '').lower() or 'blueprint' in viz.get('description', '').lower()
            
            if is_background:
                logger.info(f"Visual {i}: Detected background/blueprint. Skipping rembg.")
                crop = crop_visual_element(image, viz.get("bbox"))
            else:
                 # 標準物件 (圖示、照片) -> 使用 rembg
                crop = process_transparent_crop(image, viz.get("bbox"))
            
            if crop:
                visual_crops.append(crop)
            else:
                visual_crops.append(None) # Keep index alignment
        
        # 將裁切圖附加到分析結果 (記憶體內傳輸)
        # [v5.4 修正] 將 PIL 圖片轉換為 Base64 字串以進行 JSON 序列化
        # 這防止 "TypeError: Object of type Image is not JSON serializable"
        visual_crops_b64 = []
        for crop in visual_crops:
            if crop:
                try:
                    buf = io.BytesIO()
                    crop.save(buf, format='PNG')
                    b64_str = base64.b64encode(buf.getvalue()).decode('utf-8')
                    visual_crops_b64.append(f"data:image/png;base64,{b64_str}")
                except Exception as e:
                    logger.warning(f"Crop serialization failed: {e}")
                    visual_crops_b64.append(None)
            else:
                visual_crops_b64.append(None)

        analysis_result["_visual_crops"] = visual_crops_b64

        # [v6.1] 重建模式策略
        # 對於掃描文件/藍圖，我們偏好重建 (白板) 勝於修復 (Inpainting)。
        # 這更乾淨、更快，且避免 "灰色方塊" 偽影。
        USE_RECONSTRUCTION = True
        
        if USE_RECONSTRUCTION:
            analysis_result["reconstruction_mode"] = True
            # [v6.1 修正] 回傳帶有視覺元素的白色空白圖片
            # 這確保網頁編輯器有有效的預覽可顯示 (所見即所得)
            cleaned_image = Image.new('RGB', image.size, (255, 255, 255))
            
            # 將視覺裁切圖合成到白色背景上
            for i, crop in enumerate(visual_crops):
                if crop:
                    try:
                        bbox = visual_elements[i].get('bbox', [0,0,0,0])
                        ymin, xmin, ymax, xmax = bbox
                        
                        # 計算像素位置
                        left = int(xmin / 1000 * image.width)
                        top = int(ymin / 1000 * image.height)
                        
                        # 如果需要，調整裁切圖大小以匹配 bbox 大小 (由於裁切提取邏輯的細微差別)
                        # 通常裁切尺寸與 bbox 相符，但為了保險起見或直接貼上
                        # 目前採簡單貼上。
                        
                        # 處理透明度
                        if crop.mode in ('RGBA', 'LA'):
                             cleaned_image.paste(crop, (left, top), mask=crop)
                        else:
                             cleaned_image.paste(crop, (left, top))
                             
                    except Exception as e:
                        logger.warning(f"Failed to composite visual element {i} onto preview: {e}")
            
            logger.info(f"頁面 {page_num}: 重建模式已啟用。已建立合成的白色預覽圖片。")
        else:
            # [舊版路徑] 修補與 Inpaint
            # 2. 修補文字區域與視覺區域 (確定性遮罩)
            mask_targets = analysis_result.get('elements', []) + visual_elements
            patched_image = patch_text_areas(image, mask_targets)
            
            # 步驟 3：生成式修補 (優化背景)
            cleaned_image = await remove_text_from_image(patched_image, api_key, remove_icon, original_image=image)
            analysis_result["reconstruction_mode"] = False

        return (analysis_result, cleaned_image)
        
    except Exception as e:
        logger.error(f"Page {page_num}: Vision fallback failed: {e}")
        return ({
            "title": "處理失敗",
            "content": ["無法分析此頁面"],
            "layout": "split_left_image",
            "_visual_crops": []
        }, image)



def reconstruct_slide_background(slide, bg_hex):
    """
    [v6.1] 為重建模式生成乾淨的純色/漸層背景。
    """
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_hex)
        # 未來：根據分析添加細微的漸層或圖案
    except Exception as e:
        logger.warning(f"Background reconstruction failed: {e}")

def vectorize_image_to_svg(pil_image):
    """
    [v6.2] 使用 Potrace 將 PIL 圖片轉換為 SVG。
    回傳暫存 SVG 檔案的路徑，若失敗則回傳 None。
    """
    if not shutil.which("potrace"):
        # logger.warning("未找到 Potrace。跳過向量化。")
        return None
        
    try:
        # 建立暫存 BMP 檔案
        with tempfile.NamedTemporaryFile(suffix=".bmp", delete=False) as bmp_file:
            bmp_path = bmp_file.name
            
        svg_path = bmp_path.replace(".bmp", ".svg")
        
        # 預處理：灰階 + 閾值
        img = pil_image.convert('L')
        # 閾值：< 180 變黑 (線條)，> 180 變白 (背景)
        # 如果線條太淡，請調整此閾值
        img = img.point(lambda x: 0 if x < 200 else 255, '1')
        img.save(bmp_path)
        
        # 執行 potrace
        # -s: SVG 後端
        # --alphamax 0.2: 輕微平滑曲線
        # -k 0.5: 黑階
        # [v6.2.1] 新增逾時以防止掛起
        cmd = ["potrace", bmp_path, "-s", "-o", svg_path, "--alphamax", "0.2"]
        subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=5)
        
        # 清理 BMP
        if os.path.exists(bmp_path):
            os.unlink(bmp_path)
            
        return svg_path
    except Exception as e:
        logger.error(f"Vectorization failed: {e}")
        if os.path.exists(bmp_path):
            os.unlink(bmp_path)
        return None

def create_pptx_from_analysis(analyses: List[dict], images: List, output_path: str):
    """
    根據分析結果與原始圖片生成 PPTX 檔案 (v5.0 Overlay Layout & Legacy Split).
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    SLIDE_W_INCH = 13.333
    SLIDE_H_INCH = 7.5
    
    for i, slide_data in enumerate(analyses):
        try:
            # 建立空白投影片
            slide_layout = prs.slide_layouts[6] # 6 = Blank
            slide = prs.slides.add_slide(slide_layout)
            
            # 背景顏色 (Overlay 模式下通常被背景圖覆蓋，但保留作為底色)
            bg_hex = slide_data.get("background_color_hex", "#18181b")
            text_hex = slide_data.get("text_color_hex", "#ffffff")
            text_rgb = hex_to_rgb(text_hex)
            
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = hex_to_rgb(bg_hex)
            
            # 圖片處理 (背景圖)
            img_byte_arr = None
            cleaned_bg_img = None
            
            # 初始畫布參數 (預設全版，若有圖片則更新)
            canvas_w = SLIDE_W_INCH
            canvas_h = SLIDE_H_INCH
            canvas_left = 0.0
            canvas_top = 0.0

            if i < len(images):
                cleaned_bg_img = images[i]
                if cleaned_bg_img:
                    try:
                        # [v7.3.0] 計算智慧縮放座標 (Contain Mode)
                        img_w_px, img_h_px = cleaned_bg_img.size
                        ratio_img = img_w_px / img_h_px
                        ratio_slide = SLIDE_W_INCH / SLIDE_H_INCH
                        
                        if ratio_img > ratio_slide:
                            # 圖片較寬 (上下留白)
                            canvas_w = SLIDE_W_INCH
                            canvas_h = SLIDE_W_INCH / ratio_img
                            canvas_left = 0.0
                            canvas_top = (SLIDE_H_INCH - canvas_h) / 2
                        else:
                            # 圖片較高 (左右留白)
                            canvas_h = SLIDE_H_INCH
                            canvas_w = SLIDE_H_INCH * ratio_img
                            canvas_top = 0.0
                            canvas_left = (SLIDE_W_INCH - canvas_w) / 2

                        # [Fix v7.3.4] 強力圖像處理：處理透明度並增加 PNG 備援
                        buf = io.BytesIO()
                        try:
                            # 1. 處理透明度 (Transparency Handling)
                            # 如果是 RGBA/LA 或 P (帶透明)，合成到背景色，避免轉 RGB 變黑 (或變白框)
                            if cleaned_bg_img.mode in ('RGBA', 'LA') or (cleaned_bg_img.mode == 'P' and 'transparency' in cleaned_bg_img.info):
                                # [Fix v7.3.6] 使用投影片背景色 (bg_hex)
                                bg_rgb = tuple(hex_to_rgb(bg_hex))
                                bg_layer = Image.new('RGB', cleaned_bg_img.size, bg_rgb)
                                
                                # 必須轉為 RGBA 才能正確合成 Alpha
                                alpha_composite = cleaned_bg_img.convert('RGBA')
                                bg_layer.paste(alpha_composite, mask=alpha_composite.split()[3]) # 3=Alpha
                                cleaned_bg_img = bg_layer
                            elif cleaned_bg_img.mode != 'RGB':
                                cleaned_bg_img = cleaned_bg_img.convert('RGB')
                                
                            # 2. 嘗試儲存為 JPEG (較小)
                            cleaned_bg_img.save(buf, format='JPEG', quality=90)
                        except Exception as e_jpg:
                            logger.warning(f"Slide {i}: JPEG save failed ({e_jpg}), trying PNG fallback...")
                            try:
                                buf = io.BytesIO() # Reset buffer
                                cleaned_bg_img.save(buf, format='PNG')
                            except Exception as e_png:
                                logger.error(f"Slide {i}: Image save totally failed: {e_png}")
                                buf = None
                        
                        if buf:
                            buf.seek(0)
                            img_byte_arr = buf
                        else:
                            img_byte_arr = None
                    except Exception as e:
                        logger.error(f"Background image processing failed for slide {i}: {e}")
                        img_byte_arr = None

            # [Overlay 模式檢查]
            # layout 可能是 'overlay' (v5) 或 'split_left_image' (v2) 或其他。
            layout_type = slide_data.get('layout', 'split_left_image')
            elements = slide_data.get('elements', [])
            visual_elements = slide_data.get('visual_elements', [])
            
            # 如果有帶有 BBox 的 'elements'，強制使用 overlay 模式
            if elements and len(elements) > 0:
                layout_type = 'overlay'

            # --- 佈局實作 ---
            # --- 佈局實作 ---
            # [v6.1] 重建模式檢查
            reconstruction_mode = slide_data.get("reconstruction_mode", False)

            if layout_type == 'overlay':
                
                # 1. 背景處理
                if reconstruction_mode:
                    # [乾淨重建]
                    # 捨棄原始圖片。使用生成的乾淨背景。
                    reconstruct_slide_background(slide, bg_hex)
                    # 不為背景執行 add_picture！
                else:
                    # [舊版/復原]
                    # 使用處理過 (清理/修補) 的圖片作為背景
                    if img_byte_arr:
                        try:
                            # 使用計算出的 Canvas 座標與尺寸
                            pic = slide.shapes.add_picture(
                                img_byte_arr, 
                                Inches(canvas_left), 
                                Inches(canvas_top), 
                                width=Inches(canvas_w), 
                                height=Inches(canvas_h)
                            )
                        except Exception as e:
                            logger.error(f"Slide {i}: Add bg picture failed: {e}")
                        
                # 2. [v5.4] 視覺物件提取 (圖片/圖表)
                # 將裁切圖放置為獨立的圖片形狀
                visual_crops = slide_data.get("_visual_crops", [])
                for idx, viz in enumerate(visual_elements):
                    try:
                        if idx < len(visual_crops) and visual_crops[idx]:
                            crop_data = visual_crops[idx]
                            crop_img = None
                            
                            # 處理 Base64 字串 (v5.4 序列化修正)
                            if isinstance(crop_data, str) and crop_data.startswith("data:image"):
                                try:
                                    header, encoded = crop_data.split(",", 1)
                                    img_bytes = base64.b64decode(encoded)
                                    crop_img = Image.open(io.BytesIO(img_bytes))
                                except Exception as e:
                                    logger.warning(f"Failed to decode visual element crop: {e}")
                            # 處理 PIL 圖片 (舊版/內部)
                            elif isinstance(crop_data, Image.Image):
                                crop_img = crop_data
                                
                            if crop_img:
                                bbox = viz.get('bbox', [0,0,0,0])
                                ymin, xmin, ymax, xmax = bbox
                                
                                # [v7.3.0] 座標投影校正 (Project BBox to Scaled Background Canvas)
                                # 這是 AI 預測的「理想框」
                                target_left = Inches(canvas_left + (xmin / 1000 * canvas_w))
                                target_top = Inches(canvas_top + (ymin / 1000 * canvas_h))
                                target_width = Inches((xmax - xmin) / 1000 * canvas_w)
                                target_height = Inches((ymax - ymin) / 1000 * canvas_h)
                                
                                # [v7.3.1 Fix] 內部素材變形修正 (Smart Fit for Elements)
                                # 不強制填滿理想框，而是「保持比例放入」(Contain) 並置中
                                # 這避免了圖表被壓扁或拉長
                                
                                # 計算原始圖片比例
                                elem_w_px, elem_h_px = crop_img.size
                                elem_ratio = elem_w_px / elem_h_px
                                
                                # 計算目標框比例 (使用 Inches 值)
                                # 注意：Inches 物件可直接比較與運算
                                target_ratio = target_width / target_height
                                
                                final_width = target_width
                                final_height = target_height
                                final_left = target_left
                                final_top = target_top
                                
                                if elem_ratio > target_ratio:
                                    # 圖片較扁長 -> 寬度填滿，高度縮小，垂直置中
                                    final_width = target_width
                                    final_height = target_width / elem_ratio
                                    final_left = target_left
                                    final_top = target_top + (target_height - final_height) / 2
                                else:
                                    # 圖片較瘦高 -> 高度填滿，寬度縮小，水平置中
                                    final_height = target_height
                                    final_width = target_height * elem_ratio
                                    final_top = target_top
                                    final_left = target_left + (target_width - final_width) / 2
                                
                                # [v6.2] 向量化檢查
                                is_blueprint = 'blueprint' in viz.get('description', '').lower() or 'background' in viz.get('type', '').lower()
                                svg_path = None
                                
                                if reconstruction_mode and is_blueprint:
                                     svg_path = vectorize_image_to_svg(crop_img)
                                     
                                if svg_path:
                                     try:
                                         # 插入 SVG (向量)
                                         slide.shapes.add_picture(svg_path, final_left, final_top, width=final_width, height=final_height)
                                         # 清理
                                         if os.path.exists(svg_path):
                                             os.unlink(svg_path)
                                     except Exception as ve:
                                         logger.warning(f"SVG injection failed, falling back to raster: {ve}")
                                         svg_path = None # Trigger fallback

                                if not svg_path:
                                    # 將 PIL 轉換為位元組 (點陣圖備援)
                                    crop_buf = io.BytesIO()
                                    crop_img.save(crop_buf, format='PNG') 
                                    crop_buf.seek(0)
                                    
                                    slide.shapes.add_picture(crop_buf, final_left, final_top, width=final_width, height=final_height)
                    except Exception as e:
                        logger.warning(f"Visual element {idx} placement failed: {e}")

                # 3. 文字覆蓋
                for elem in elements:
                    try:
                        content = elem.get('content', '')
                        bbox = elem.get('bbox', [0,0,0,0]) # [ymin, xmin, ymax, xmax] 0-1000
                        font_size_pt = elem.get('font_size', 24)
                        align = elem.get('alignment', 'left')
                        color_hex = elem.get('color_hex', text_hex)
                        
                        ymin, xmin, ymax, xmax = bbox
                        
                        # [v7.3.0] 座標投影校正
                        left = Inches(canvas_left + (xmin / 1000 * canvas_w))
                        top = Inches(canvas_top + (ymin / 1000 * canvas_h))
                        width = Inches((xmax - xmin) / 1000 * canvas_w)
                        height = Inches((ymax - ymin) / 1000 * canvas_h)
                        
                        # 新增文字方塊
                        tb = slide.shapes.add_textbox(left, top, width, height)
                        tf = tb.text_frame
                        tf.word_wrap = True
                        
                        p = tf.paragraphs[0]
                        p.text = str(content)
                        
                        # 如果字體大小缺失或太小，使用啟發式方法
                        if not isinstance(font_size_pt, (int, float)) or font_size_pt < 8: 
                            font_size_pt = 18
                        
                        p.font.size = Pt(font_size_pt)
                        try:
                            p.font.color.rgb = hex_to_rgb(color_hex)
                        except:
                            p.font.color.rgb = text_rgb
                        
                        if align == 'center':
                            p.alignment = PP_ALIGN.CENTER
                        elif align == 'right':
                            p.alignment = PP_ALIGN.RIGHT
                        else:
                            p.alignment = PP_ALIGN.LEFT
                            
                    except Exception as e:
                         # logger.warning(f"Element rendering failed: {e}")
                         pass

            elif layout_type == 'full_width_text':
                 if slide_data.get("title"):
                    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.3), Inches(1.5))
                    title_tf = title_box.text_frame
                    title_tf.word_wrap = True
                    title_p = title_tf.paragraphs[0]
                    title_p.text = slide_data["title"]
                    title_p.font.size = Pt(36)
                    title_p.font.bold = True
                    title_p.font.color.rgb = text_rgb
                    title_p.alignment = PP_ALIGN.CENTER
                
                 content_items = slide_data.get("content", [])
                 if content_items:
                    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(4.5))
                    content_tf = content_box.text_frame
                    content_tf.word_wrap = True
                    for item in content_items:
                        p = content_tf.add_paragraph()
                        p.text = str(item)
                        p.font.size = Pt(20)
                        p.font.color.rgb = text_rgb
                        p.space_after = Pt(20)
            else:
                # 分割佈局 (v2 舊版/備援預設)
                if img_byte_arr:
                    # 左側圖片
                    try:
                        pic = slide.shapes.add_picture(img_byte_arr, Inches(0), Inches(0), height=prs.slide_height)
                        # 如果需要，置中於左半部 (Inches(6.6))，但全高適合分割佈局
                        # 如果太寬則裁切
                        if pic.width > Inches(7):
                             crop = (pic.width - Inches(7)) / 2
                             pic.crop_left = crop / pic.width
                             pic.crop_right = crop / pic.width
                             pic.left = 0
                    except Exception as e:
                        logger.error(f"Slide {i}: Add picture failed: {e}")

                # 右側文字
                text_left = Inches(7.0)
                text_width = Inches(5.8)
                
                if slide_data.get("title"):
                    title_box = slide.shapes.add_textbox(text_left, Inches(0.5), text_width, Inches(1.5))
                    title_tf = title_box.text_frame
                    title_tf.word_wrap = True
                    title_p = title_tf.paragraphs[0]
                    title_p.text = slide_data["title"]
                    title_p.font.size = Pt(28)
                    title_p.font.bold = True
                    title_p.font.color.rgb = text_rgb

                content_items = slide_data.get("content", [])
                if content_items:
                    content_box = slide.shapes.add_textbox(text_left, Inches(2.2), text_width, Inches(4.5))
                    content_tf = content_box.text_frame
                    content_tf.word_wrap = True
                    for item in content_items:
                        p = content_tf.add_paragraph()
                        p.text = str(item)
                        p.font.size = Pt(16)
                        p.font.color.rgb = text_rgb
                        p.space_after = Pt(12)
                        p.level = 0
            
            # 演講者筆記
            if slide_data.get("speaker_notes"):
                slide.notes_slide.notes_text_frame.text = slide_data["speaker_notes"]

        except Exception as e:
            logger.error(f"Slide {i} generation failed: {e}")
            continue

    prs.save(output_path)
    logger.info(f"簡報已儲存至: {output_path}")

def generate_preview_images(pdf_bytes: bytes, output_dir: str) -> List[str]:
    """
    [v7.1] 生成 PDF 預覽縮圖，回傳 Base64 Data URL (Stateless)
    修正 Cloud Run 上因無狀態導致 /static/temp/ 404 的問題
    """
    try:
        # 降低記憶體使用量：dpi=150 (縮圖可較低)，thread_count=1
        images = convert_from_bytes(pdf_bytes, dpi=150, thread_count=1)
        logger.info(f"預覽生成: 轉換了 {len(images)} 張圖片")
        
        image_data_urls = []
        for i, img in enumerate(images):
            try:
                # 調整為縮圖大小
                img.thumbnail((400, 400))
                
                # 如果需要，轉換為 RGB (JPEG 不支援透明度)
                if img.mode in ('RGBA', 'LA', 'P'):
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    if img.mode in ('RGBA', 'LA'):
                        background.paste(img, mask=img.split()[-1])
                    img = background
                elif img.mode != 'RGB':
                    img = img.convert('RGB')
                
                # 轉換為 Base64
                buf = io.BytesIO()
                img.save(buf, format='JPEG', quality=70, optimize=True)
                b64_str = base64.b64encode(buf.getvalue()).decode('utf-8')
                image_data_urls.append(f"data:image/jpeg;base64,{b64_str}")
                
            except Exception as page_err:
                logger.warning(f"頁面 {i+1} 預覽生成失敗: {page_err}")
                # 使用錯誤佔位符
                image_data_urls.append("data:image/svg+xml;base64,PHN2ZyB4bWxucz0iaHR0cDovL3d3dy53My5vcmcvMjAwMC9zdmciIHZpZXdCb3g9IjAgMCAxMDAgMTAwIj48cmVjdCB3aWR0aD0iMTAwIiBoZWlnaHQ9IjEwMCIgZmlsbD0iIzMzMyIgLz48dGV4dCB4PSI1MCIgeT0iNTUiIGZpbGw9IiNhYWEiIGZvbnQtc2l6ZT0iMTAiIHRleHQtYW5jaG9yPSJtaWRkbGUiPuWksei0qzwvdGV4dD48L3N2Zz4=")
        
        return image_data_urls
    except Exception as e:
        logger.error(f"預覽生成失敗 (Memory/Poppler): {e}")
        raise ValueError(f"無法生成預覽: {str(e)}")


from pdf2image import convert_from_path
from pypdf import PdfReader

async def analyze_presentation(pdf_path: str, api_key: str, filename: str, selected_indices: Optional[List[int]] = None, remove_icon: bool = False, progress_callback: Optional[callable] = None) -> tuple:
    """
    主要流程：PDF (File) -> 圖片 -> Gemini 分析 -> 文字移除
    優化 (v2.10.17): pypdf 秒讀頁數 + 首頁優先策略 (Priority First Page) + 強制超時保護。
    """
    logger.info(f"開始處理 PDF: {filename} (Path: {pdf_path})")
    
    # [最佳化] 立即通知開始，以更新 UI 狀態 "準備中"
    if progress_callback:
        try:
             await progress_callback(0, 0, message="正在讀取 PDF 檔案結構...")
        except:
             pass

    try:
        # 1. 快速獲取 PDF 資訊 (使用 pypdf) - 強制 10s Timeout
        # 在線程中執行，因為 pypdf 檔案讀取是同步 IO
        def get_pdf_count():
            with open(pdf_path, 'rb') as f:
                reader = PdfReader(f)
                return len(reader.pages)
        
        # PDF 讀取逾時保護 (10秒)
        total_pdf_pages = await asyncio.wait_for(
            asyncio.to_thread(get_pdf_count), 
            timeout=10
        )
        logger.info(f"PDF 總頁數: {total_pdf_pages}")
        
    except asyncio.TimeoutError:
        logger.error("PDF 讀取超時 (pypdf)")
        raise ValueError("PDF 檔案讀取超時，請檢查檔案是否損毀或過大")
        
    except Exception as e:
        logger.error(f"無法讀取 PDF 資訊: {e}")
        raise ValueError(f"無法讀取 PDF 結構: {str(e)}")

    # 決定要處理的頁面索引 (0-based)
    target_indices = selected_indices if selected_indices else list(range(total_pdf_pages))
    target_indices = [i for i in target_indices if 0 <= i < total_pdf_pages]
    target_indices.sort()
    
    if not target_indices:
        return [], []

    analyses = []
    cleaned_images = []
    
    # 策略配置
    BATCH_SIZE = 3
    DELAY_BETWEEN_BATCHES = 1
    TIMEOUT_PER_BATCH = 45 # seconds (Image Conversion)
    TIMEOUT_PER_PAGE_ANALYSIS = 90 # seconds (Extended for high-res)
    


    total_target = len(target_indices)
    
    # "優先首頁" 的自訂迴圈
    # 我們手動建構批次以確保第 1 批次為單頁 (為了速度)
    batches = []
    remaining_indices = target_indices.copy()
    
    # 設定第一批次 (優先)
    if remaining_indices:
        # 第一批次只有 1 頁以確保即時回饋
        batches.append([remaining_indices.pop(0)])
    
    # 設定後續批次
    while remaining_indices:
        chunk = remaining_indices[:BATCH_SIZE]
        batches.append(chunk)
        remaining_indices = remaining_indices[BATCH_SIZE:]

    # 執行批次
    processed_count = 0
    
    for batch_indices in batches:
        current_batch_size = len(batch_indices)
        
        # 通知進度：轉檔中
        if progress_callback:
            start_p = batch_indices[0] + 1
            end_p = batch_indices[-1] + 1
            msg = f"正在處理第 {start_p}/{total_target} 頁..."
            try:
                await progress_callback(processed_count, total_target, message=msg)
            except Exception as e:
                logger.warning(f"Callback msg failed: {e}")

        # 2. 隨需影像轉換 (受逾時保護)
        batch_images = []
        try:
            # 檢查連續性以進行最佳化
            is_consecutive = (len(batch_indices) > 1 and 
                             batch_indices[-1] - batch_indices[0] == len(batch_indices) - 1)
            
            async def run_conversion():
                if is_consecutive:
                    # 範圍轉換的相異參數
                    s = batch_indices[0] + 1
                    e = batch_indices[-1] + 1
                    return await asyncio.to_thread(
                        convert_from_path, pdf_path, dpi=200, 
                        first_page=s, last_page=e, thread_count=1
                    )
                else:
                    imgs = []
                    for idx in batch_indices:
                        p = idx + 1
                        res = await asyncio.to_thread(
                           convert_from_path, pdf_path, dpi=200,
                           first_page=p, last_page=p, thread_count=1
                        )
                        if res: imgs.extend(res)
                    return imgs

            # 逾時包裝器
            batch_images = await asyncio.wait_for(run_conversion(), timeout=TIMEOUT_PER_BATCH)
            
        except asyncio.TimeoutError:
            logger.error(f"Batch conversion timed out after {TIMEOUT_PER_BATCH}s")
            batch_images = [Image.new('RGB', (800, 600), color='white') for _ in batch_indices]
        except Exception as e:
            logger.error(f"Batch conversion failed: {e}")
            batch_images = [Image.new('RGB', (800, 600), color='white') for _ in batch_indices]

        # 安全填充
        while len(batch_images) < current_batch_size:
             batch_images.append(Image.new('RGB', (800, 600), color='white'))
        batch_images = batch_images[:current_batch_size]

        # 3. 分析批次
        tasks = [
            process_single_page(img, batch_indices[j] + 1, total_pdf_pages, api_key, remove_icon=remove_icon, pdf_path=pdf_path)
            for j, img in enumerate(batch_images)
        ]
        
        batch_results = await asyncio.gather(*tasks, return_exceptions=True)
        
        # 立即收集結果
        for res in batch_results:
            if isinstance(res, Exception):
                logger.error(f"Batch Analysis Critical Failure: {res}")
                analyses.append({
                    "title": "分析失敗", 
                    "content": ["請手動編輯此頁面"], 
                    "layout": "split_left_image"
                })
                cleaned_images.append(Image.new('RGB', (800, 600), color='white')) 
            elif isinstance(res, tuple) and len(res) == 2:
                # 確保 analysis 是 dict
                analysis = res[0]
                if not isinstance(analysis, dict):
                    analysis = {"title": "格式錯誤", "content": [str(analysis)[:200]]}
                analyses.append(analysis)
                cleaned_images.append(res[1])
            else:
                 # 未預期的格式，給預設值
                 logger.warning(f"Unexpected result format: {res}")
                 analyses.append({"title": "未知錯誤", "content": []})
                 cleaned_images.append(Image.new('RGB', (800, 600), color='white'))

        processed_count += current_batch_size
        
        # [安全策略] 首頁探測
        # 如果第一頁分析失敗，則中止整個流程。
        # 這可以防止在 PDF 格式根本無法讀取或 API 停機的情況下浪費 50 多頁的成本。
        if processed_count == 1 and len(analyses) > 0:
            first_result = analyses[0]
            # 檢查已知的錯誤簽名
            error_titles = ["Analysis Error", "Parse Error", "分析失敗", "未知錯誤", "處理失敗"]
            if (first_result.get("title") in error_titles) or (not first_result.get("content")):
                logger.error("🛑 首頁探測失敗。中止剩餘頁面以節省成本。")
                # 我們在此停止。UI 將只收到這一個錯誤投影片。
                break

        if processed_count < total_target:
            await asyncio.sleep(DELAY_BETWEEN_BATCHES)

    logger.info(f"所有頁面處理完成。總共: {len(analyses)} 頁")
        
    return analyses, cleaned_images


async def process_pdf_to_slides(pdf_content, api_key: str, filename: str, selected_indices: List[int] = None):
    # 舊版包裝器
    analyses, cleaned_images = await analyze_presentation(pdf_content, api_key, filename, selected_indices)
    output_dir = "temp_slides"
    os.makedirs(output_dir, exist_ok=True)
    base_name = os.path.splitext(filename)[0]
    output_path = os.path.join(output_dir, f"{base_name}_converted.pptx")
    await asyncio.to_thread(create_pptx_from_analysis, analyses, cleaned_images, output_path)
    return output_path

