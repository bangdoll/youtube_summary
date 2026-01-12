
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def analyze_pptx(path):
    print(f"--- Analyzing: {os.path.basename(path)} ---")
    try:
        prs = Presentation(path)
        print(f"Total Slides: {len(prs.slides)}")
        print(f"Slide Size: {prs.slide_width.inches} x {prs.slide_height.inches} inches")
        
        for i, slide in enumerate(prs.slides[:3]): # Analyze first 3 slides
            print(f"\nSlide {i+1}:")
            # Background
            bg = slide.background
            try:
                if bg.fill.type:
                    print(f"  Background: {bg.fill.type}") # SOLID, GRADIENT, PICTURE etc
            except:
                print("  Background: Default/None")

            # Shapes
            text_boxes = 0
            pictures = 0
            shapes = 0
            groups = 0
            
            for shape in slide.shapes:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        text_boxes += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        pictures += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        # Check if it contains text
                        if shape.has_text_frame and shape.text.strip():
                            text_boxes += 1
                        else:
                            shapes += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        groups += 1
                    else:
                        shapes += 1
                    
                    # Text Content Sample
                    if shape.has_text_frame and shape.text.strip():
                        print(f"  > Text Found ({shape.shape_type}): {shape.text[:50]}...")
                except Exception as shape_error:
                    # Skip confusing shapes
                    continue

    except Exception as e:
        print(f"Error analyzing {path}: {e}")

if __name__ == "__main__":
    analyze_pptx("Docs/Awakening_Blueprint.pptx")
    print("\n" + "="*30 + "\n")
    analyze_pptx("Docs/Awakening_Blueprint (7).pptx")
